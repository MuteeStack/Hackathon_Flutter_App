"""
Generates a premium supporting slide deck (PPTX) for the QuickFix solution.
Uses python-pptx to build slides with premium design colors:
Dark Navy (0x0D, 0x1E, 0x3D), Brand Pink (0xEA, 0x4C, 0x6A), and Gemini Purple (0x9B, 0x72, 0xCB).
"""
import sys
import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    import subprocess
    print("python-pptx not found. Installing python-pptx...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    
    # 16:9 widescreen layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_slide_layout = prs.slide_layouts[6]
    
    # Premium Color Palette
    DARK_NAVY = RGBColor(13, 30, 61)       # #0D1E3D
    BRAND_PINK = RGBColor(234, 76, 106)     # #EA4C6A
    GEMINI_PURPLE = RGBColor(155, 114, 203) # #9B72CB
    GEMINI_BLUE = RGBColor(66, 133, 244)    # #4285F4
    WHITE = RGBColor(255, 255, 255)
    LIGHT_GRAY = RGBColor(245, 246, 249)
    DARK_TEXT = RGBColor(31, 31, 31)
    MUTED_TEXT = RGBColor(95, 99, 104)
    
    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_header(slide, title_text, category_text="QUICKFIX SOLUTION DECK", is_dark=False):
        # Category Tracker
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
        cat_tf = cat_box.text_frame
        cat_tf.word_wrap = True
        cat_tf.margin_left = cat_tf.margin_top = cat_tf.margin_bottom = cat_tf.margin_right = 0
        p_cat = cat_tf.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.name = 'Arial'
        p_cat.font.color.rgb = BRAND_PINK if is_dark else GEMINI_PURPLE
        
        # Main Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_tf.margin_left = title_tf.margin_top = title_tf.margin_bottom = title_tf.margin_right = 0
        p_title = title_tf.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(28)
        p_title.font.bold = True
        p_title.font.name = 'Arial'
        p_title.font.color.rgb = WHITE if is_dark else DARK_NAVY

    # ============================================================
    # SLIDE 1: Title Slide (Dark Background)
    # ============================================================
    slide1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide1, DARK_NAVY)
    
    # Premium decorative shape (Gradient-like effect using overlapping panels)
    shape = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BRAND_PINK
    shape.line.fill.background()
    
    shape2 = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0), Inches(0.2), Inches(7.5)
    )
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = GEMINI_PURPLE
    shape2.line.fill.background()

    # Title Text Frame
    tb = slide1.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10), Inches(3.5))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "QUICKFIX"
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Arial'
    
    p2 = tf.add_paragraph()
    p2.text = "Hybrid Agentic AI Service Orchestrator"
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = BRAND_PINK
    p2.font.name = 'Arial'
    p2.space_before = Pt(10)
    
    p3 = tf.add_paragraph()
    p3.text = "Challenge 2: AI Service Orchestrator for Pakistan's Informal Economy\nPowered by Google Antigravity & Gemini 2.0 Flash"
    p3.font.size = Pt(16)
    p3.font.color.rgb = GEMINI_PURPLE
    p3.font.name = 'Arial'
    p3.space_before = Pt(30)

    # ============================================================
    # SLIDE 2: The Problem (Light Background)
    # ============================================================
    slide2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide2, WHITE)
    add_header(slide2, "Pakistan's Informal Economy Challenges")
    
    # Left Content Column (Text bullet points)
    tb_left = slide2.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5))
    tf_left = tb_left.text_frame
    tf_left.word_wrap = True
    
    p_header = tf_left.paragraphs[0]
    p_header.text = "The Status Quo: Fragmented & Manual"
    p_header.font.size = Pt(20)
    p_header.font.bold = True
    p_header.font.color.rgb = DARK_NAVY
    p_header.space_after = Pt(15)
    
    bullets = [
        "● WhatsApp & Phone Calls: Service matching happens in chaotic, unstructured chats with zero central coordination.",
        "● High Disruption & Missed Work: Gig-workers miss customer opportunities due to lack of automation or scheduling systems.",
        "● Manual Referrals: Customers struggle to find nearby, trusted, and verified professionals in real-time.",
        "● Lack of Uptime & Connectivity: Systems relying 100% on heavy LLM cloud APIs fail in areas with spotty network coverage."
    ]
    for bullet in bullets:
        p_b = tf_left.add_paragraph()
        p_b.text = bullet
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = DARK_TEXT
        p_b.space_after = Pt(12)
        p_b.font.name = 'Arial'

    # Right Content Column (Visual Summary Cards)
    card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.5))
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT_GRAY
    card.line.color.rgb = GEMINI_PURPLE
    card.line.width = Pt(1.5)
    
    tf_card = card.text_frame
    tf_card.word_wrap = True
    tf_card.margin_left = tf_card.margin_top = tf_card.margin_bottom = tf_card.margin_right = Inches(0.4)
    
    p_c1 = tf_card.paragraphs[0]
    p_c1.text = "Why Traditional Apps Fail here:"
    p_c1.font.size = Pt(18)
    p_c1.font.bold = True
    p_c1.font.color.rgb = DARK_NAVY
    p_c1.space_after = Pt(20)
    
    sub_points = [
        "1. Complex Interfaces: Gig-workers cannot easily navigate complex forms.",
        "2. Strict Language Barriers: Pakistan's informal workers rely heavily on Roman Urdu ('kal subah AC theek krna hai') and local scripts which traditional apps fail to parse.",
        "3. Network Fragility: Extreme dependency on continuous high-speed internet causes total application lockouts."
    ]
    for sp in sub_points:
        p_sp = tf_card.add_paragraph()
        p_sp.text = sp
        p_sp.font.size = Pt(13)
        p_sp.font.color.rgb = MUTED_TEXT
        p_sp.space_after = Pt(15)

    # ============================================================
    # SLIDE 3: The Solution (Light Background)
    # ============================================================
    slide3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide3, WHITE)
    add_header(slide3, "The Solution: QuickFix Agentic AI Service Platform")

    # Full Width Text Box
    tb_sol = slide3.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf_sol = tb_sol.text_frame
    tf_sol.word_wrap = True
    
    p_intro = tf_sol.paragraphs[0]
    p_intro.text = "QuickFix automates the end-to-end booking lifecycle, powered by Google Antigravity Orchestration and a unique Hybrid Approach."
    p_intro.font.size = Pt(16)
    p_intro.font.italic = True
    p_intro.font.color.rgb = GEMINI_PURPLE
    p_intro.space_after = Pt(24)

    # 3 Column features
    col_width = Inches(3.6)
    col_gap = Inches(0.4)
    col_top = Inches(2.8)
    col_height = Inches(3.8)
    
    columns_data = [
        {
            "title": "🗣 Intent Understanding",
            "desc": "Accepts speech or text inputs in English, Urdu, and Roman Urdu. Parses complex, informal user requests into structured intents seamlessly.",
            "color": GEMINI_BLUE
        },
        {
            "title": "📍 Proactive Discover & Match",
            "desc": "Scans Mock Datasets and real Google Maps Places. Ranks matches using Gemini AI's cognitive scores, providing transparent reasoning for its choices.",
            "color": BRAND_PINK
        },
        {
            "title": "⚡ Simulation & Follow-up",
            "desc": "Simulates secure booking transactions, updates cloud Firestore/local databases, generates PDF receipts, and schedules proactive reminders.",
            "color": GEMINI_PURPLE
        }
    ]
    
    for i, col in enumerate(columns_data):
        left_pos = Inches(0.8) + i * (col_width + col_gap)
        c_shape = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, col_top, col_width, col_height)
        c_shape.fill.solid()
        c_shape.fill.fore_color.rgb = LIGHT_GRAY
        c_shape.line.color.rgb = col["color"]
        c_shape.line.width = Pt(2)
        
        c_tf = c_shape.text_frame
        c_tf.word_wrap = True
        c_tf.margin_left = c_tf.margin_top = c_tf.margin_bottom = c_tf.margin_right = Inches(0.3)
        
        cp_title = c_tf.paragraphs[0]
        cp_title.text = col["title"]
        cp_title.font.size = Pt(16)
        cp_title.font.bold = True
        cp_title.font.color.rgb = DARK_NAVY
        cp_title.space_after = Pt(15)
        
        cp_desc = c_tf.add_paragraph()
        cp_desc.text = col["desc"]
        cp_desc.font.size = Pt(12)
        cp_desc.font.color.rgb = DARK_TEXT
        cp_desc.space_after = Pt(10)

    # ============================================================
    # SLIDE 4: Core Innovation: "Hybrid Resilient Model" (Light BG)
    # ============================================================
    slide4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide4, WHITE)
    add_header(slide4, "The Core Innovation: 5-Pillar Hybrid Model")
    
    tb_innov = slide4.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf_innov = tb_innov.text_frame
    tf_innov.word_wrap = True
    
    innovations = [
        ("1. Hybrid NLU Parsing (Gemini AI + Local Heuristics)", 
         "Uses Gemini 2.0 Flash to comprehend language scripts. If internet connection is spotty, it activates an offline keyword parser mapping common Urdu/Roman words."),
        ("2. Hybrid Discovery (Google Maps Places + Sector Mock DB)", 
         "Combines live coordinates from Google Geocoding/Places with a detailed localized fallback database for offline sector positioning."),
        ("3. Hybrid Decision Matching (Generative Rating + Proximity sorting)", 
         "Uses Gemini as a Matching Expert to weigh customer urgency, ratings, and experience. If rate-limited, it implements a proximity-based Haversine sort."),
        ("4. Hybrid Cross-Platform Runs (Android APK + Web browser)", 
         "The app compiles to a native Android APK using tree-shaking and SDK 35, while fully rendering inside Google Chrome for easy testing."),
        ("5. Hybrid Storage Sync (Firebase Firestore + Local JSON)", 
         "Syncs active states to Firebase real-time nodes. Automatically writes to local JSON fallback configurations for full standalone server execution.")
    ]
    
    for title, desc in innovations:
        p_t = tf_innov.add_paragraph()
        p_t.text = title
        p_t.font.size = Pt(15)
        p_t.font.bold = True
        p_t.font.color.rgb = BRAND_PINK
        
        p_d = tf_innov.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = DARK_TEXT
        p_d.space_after = Pt(14)

    # ============================================================
    # SLIDE 5: Multi-Agent Orchestrator Pipeline (Dark Background)
    # ============================================================
    slide5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide5, DARK_NAVY)
    add_header(slide5, "The Google Antigravity Multi-Agent Pipeline", "ARCHITECTURE LAYER", is_dark=True)
    
    # Draw pipeline cards programmatically
    c_width = Inches(2.1)
    c_height = Inches(4.2)
    c_gap = Inches(0.2)
    c_top = Inches(2.2)
    
    agents = [
        {"num": "1", "name": "Intent Agent", "role": "NLU Parsing", "detail": "Extracts service, location, urgency from Roman Urdu using Gemini & regex fallback.", "color": GEMINI_BLUE},
        {"num": "2", "name": "Discovery Agent", "role": "Geo-matching", "detail": "Queries mock database and Google Places to find all available services near user's sector.", "color": GEMINI_PURPLE},
        {"num": "3", "name": "Matching Agent", "role": "Cognitive Score", "detail": "Acts as the expert to score (0-100) and rank providers with detailed customer reasoning.", "color": BRAND_PINK},
        {"num": "4", "name": "Booking Agent", "role": "Action Simulation", "detail": "Generates unique booking receipts, commits to Firestore, and drafts WhatsApp templates.", "color": GEMINI_BLUE},
        {"num": "5", "name": "Follow-Up Agent", "role": "Job Reminders", "detail": "Triggers background listeners. Schedules a proactive reminder 60m before job time.", "color": GEMINI_PURPLE}
    ]
    
    for idx, agent in enumerate(agents):
        x_pos = Inches(0.8) + idx * (c_width + c_gap)
        c_shape = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, c_top, c_width, c_height)
        c_shape.fill.solid()
        c_shape.fill.fore_color.rgb = RGBColor(22, 40, 75) # A lighter navy
        c_shape.line.color.rgb = agent["color"]
        c_shape.line.width = Pt(1.5)
        
        c_tf = c_shape.text_frame
        c_tf.word_wrap = True
        c_tf.margin_left = c_tf.margin_top = c_tf.margin_bottom = c_tf.margin_right = Inches(0.2)
        
        p_n = c_tf.paragraphs[0]
        p_n.text = f"Agent {agent['num']}"
        p_n.font.size = Pt(12)
        p_n.font.bold = True
        p_n.font.color.rgb = agent["color"]
        p_n.space_after = Pt(5)
        
        p_na = c_tf.add_paragraph()
        p_na.text = agent["name"]
        p_na.font.size = Pt(15)
        p_na.font.bold = True
        p_na.font.color.rgb = WHITE
        p_na.space_after = Pt(2)
        
        p_ro = c_tf.add_paragraph()
        p_ro.text = agent["role"].upper()
        p_ro.font.size = Pt(9)
        p_ro.font.bold = True
        p_ro.font.color.rgb = GEMINI_PURPLE
        p_ro.space_after = Pt(15)
        
        p_det = c_tf.add_paragraph()
        p_det.text = agent["detail"]
        p_det.font.size = Pt(11)
        p_det.font.color.rgb = LIGHT_GRAY

    # ============================================================
    # SLIDE 6: Action Simulation & Follow-Up Automation (Light BG)
    # ============================================================
    slide6 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide6, WHITE)
    add_header(slide6, "Resilient Action Simulation & Proactive Flow")
    
    # Left Box (Simulated Booking receipt details)
    box_left = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5))
    box_left.fill.solid()
    box_left.fill.fore_color.rgb = LIGHT_GRAY
    box_left.line.color.rgb = BRAND_PINK
    box_left.line.width = Pt(1)
    
    tf_left = box_left.text_frame
    tf_left.word_wrap = True
    tf_left.margin_left = tf_left.margin_top = tf_left.margin_bottom = tf_left.margin_right = Inches(0.4)
    
    p_b1 = tf_left.paragraphs[0]
    p_b1.text = "Simulated Booking Actions"
    p_b1.font.size = Pt(18)
    p_b1.font.bold = True
    p_b1.font.color.rgb = DARK_NAVY
    p_b1.space_after = Pt(15)
    
    sim_points = [
        "● Booking ID Generation: Creates transaction receipts keyed as `BK-XXXXXX`.",
        "● Auto-Urgency Bookings: Bypasses user selection when urgency is high, ensuring rapid match.",
        "● Dual Data Storage: Updates cloud Firestore nodes and writes locally to the `bookings.json` fallback database.",
        "● Professional Receipts: Auto-generates clean text receipt payloads with details, provider contacts, and time slots."
    ]
    for point in sim_points:
        p_p = tf_left.add_paragraph()
        p_p.text = point
        p_p.font.size = Pt(12)
        p_p.font.color.rgb = DARK_TEXT
        p_p.space_after = Pt(12)

    # Right Box (Follow-up loops)
    box_right = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.5))
    box_right.fill.solid()
    box_right.fill.fore_color.rgb = LIGHT_GRAY
    box_right.line.color.rgb = GEMINI_BLUE
    box_right.line.width = Pt(1)
    
    tf_right = box_right.text_frame
    tf_right.word_wrap = True
    tf_right.margin_left = tf_right.margin_top = tf_right.margin_bottom = tf_right.margin_right = Inches(0.4)
    
    p_r1 = tf_right.paragraphs[0]
    p_r1.text = "Automated Follow-up Loops"
    p_r1.font.size = Pt(18)
    p_r1.font.bold = True
    p_r1.font.color.rgb = DARK_NAVY
    p_r1.space_after = Pt(15)
    
    f_points = [
        "● Proactive Reminders: Dynamically calculates the appointment offset to send alerts 60 minutes beforehand.",
        "● Real-time WebSockets: Supports live uvicorn WebSocket triggers for user client updates.",
        "● Feedback Surveys: Automatically stages customer satisfaction prompts 30 minutes after completion.",
        "● Resilient Status Checks: Evaluates active background threads for seamless job reminders."
    ]
    for f_pt in f_points:
        p_fp = tf_right.add_paragraph()
        p_fp.text = f_pt
        p_fp.font.size = Pt(12)
        p_fp.font.color.rgb = DARK_TEXT
        p_fp.space_after = Pt(12)

    # ============================================================
    # SLIDE 7: Summary and Judge Walkthrough UI (Light BG)
    # ============================================================
    slide7 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide7, WHITE)
    add_header(slide7, "Judge Walkthrough & Visual Trace Screen")

    tb_walk = slide7.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf_walk = tb_walk.text_frame
    tf_walk.word_wrap = True
    
    p_w1 = tf_walk.paragraphs[0]
    p_w1.text = "Why QuickFix stands out to Hackathon Judges:"
    p_w1.font.size = Pt(18)
    p_w1.font.bold = True
    p_w1.font.color.rgb = DARK_NAVY
    p_w1.space_after = Pt(15)
    
    reasons = [
        ("🎥 Live Trace Debugging Tab:", "A beautiful, transparent terminal trace built directly into the Flutter app layout. It lets judges inspect every tool execution, API fetch, intent analysis, and agentic loop in real-time."),
        ("🛡️ Extreme Resiliency:", "Demonstrates robust handling of Pakistani contexts: works flawlessly online or offline using hybrid fallback architecture."),
        ("⚙️ Autonomy & Execution:", "Google Antigravity coordinates intent, mapping, scoring, matching, action writes, and scheduling. It is not just a UI mockup, but a working end-to-end backend orchestration engine."),
        ("📦 Ready Android Build:", "Compiled Release APK (SDK 35) is pre-built, tree-shaken, and optimized for immediate evaluation on physical devices alongside Google Chrome support.")
    ]
    
    for r_title, r_desc in reasons:
        p_rt = tf_walk.add_paragraph()
        p_rt.text = r_title
        p_rt.font.size = Pt(14)
        p_rt.font.bold = True
        p_rt.font.color.rgb = BRAND_PINK
        
        p_rd = tf_walk.add_paragraph()
        p_rd.text = r_desc
        p_rd.font.size = Pt(12)
        p_rd.font.color.rgb = DARK_TEXT
        p_rd.space_after = Pt(10)

    # Save presentation
    output_path = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "QuickFix_Solution_Deck.pptx")
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")
    return output_path

if __name__ == "__main__":
    create_deck()
